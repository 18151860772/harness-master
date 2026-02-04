"""
动量守恒定律PPT生成脚本
使用python-pptx库创建完整的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 定义配色方案
BLUE = RGBColor(30, 100, 180)      # 深蓝色（标题）
ORANGE = RGBColor(230, 140, 50)    # 橙色（强调）
BLACK = RGBColor(0, 0, 0)          # 黑色（正文）
WHITE = RGBColor(255, 255, 255)    # 白色
GRAY = RGBColor(100, 100, 100)     # 灰色

def create_title_slide(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "动量守恒定律"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "初中物理"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

def create_toc_slide(prs):
    """第2页：目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 目录内容
    toc_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    toc_frame = toc_box.text_frame
    toc_frame.word_wrap = True

    items = [
        "一、动量的概念",
        "二、动量守恒定律",
        "三、动量守恒定律的应用",
        "四、典型例题分析"
    ]

    for item in items:
        p = toc_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(28)
        p.font.color.rgb = BLACK
        p.space_before = Pt(18)
        p.level = 0

def create_concept_slide_1(prs):
    """第3页：动量的定义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "一、动量的概念"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容框
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    # 定义
    p = content_frame.add_paragraph()
    p.text = "1. 定义"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "物体的质量和速度的乘积叫做动量"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(24)
    p.level = 1

    # 公式
    p = content_frame.add_paragraph()
    p.text = "2. 公式"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "p = mv"
    p.font.size = Pt(32)
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "其中：p - 动量，m - 质量，v - 速度"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.level = 1

def create_concept_slide_2(prs):
    """第4页：动量的单位"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "动量的单位与方向"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    # 单位
    p = content_frame.add_paragraph()
    p.text = "3. 单位"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "千克·米/秒 (kg·m/s)"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(24)
    p.level = 1

    # 方向
    p = content_frame.add_paragraph()
    p.text = "4. 方向"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "动量是矢量，方向与速度方向相同"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(12)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "注：动量是状态量"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.level = 1

def create_concept_slide_3(prs):
    """第5页：动量的变化"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "动量的变化"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "动量的变化量："
    p.font.size = Pt(28)
    p.font.color.rgb = BLACK
    p.space_after = Pt(18)

    p = content_frame.add_paragraph()
    p.text = "Δp = p₂ - p₁ = mv₂ - mv₁ = m(v₂ - v₁) = mΔv"
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.space_after = Pt(30)

    p = content_frame.add_paragraph()
    p.text = "说明："
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    points = [
        "动量的变化量也是矢量",
        "运算遵循平行四边形定则",
        "在同一直线上时，可规定正方向"
    ]

    for point in points:
        p = content_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
        p.level = 1

def create_law_slide_1(prs):
    """第6页：动量守恒定律内容"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "二、动量守恒定律"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "内容："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "如果一个系统不受外力，或者所受外力的矢量和为零，"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "这个系统的总动量保持不变。"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(24)

    p = content_frame.add_paragraph()
    p.text = "公式："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "p₁ + p₂ = p₁' + p₂'"
    p.font.size = Pt(32)
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "或 m₁v₁ + m₂v₂ = m₁v₁' + m₂v₂'"
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.level = 1

def create_law_slide_2(prs):
    """第7页：守恒条件"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "动量守恒的条件"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
    content_frame = content_box.text_frame

    conditions = [
        "系统不受外力或所受外力的矢量和为零",
        "系统内力远大于外力（作用时间极短）",
        "系统在某一方向上不受外力或合外力为零"
    ]

    for i, condition in enumerate(conditions, 1):
        p = content_frame.add_paragraph()
        p.text = f"{i}. {condition}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(20)
        p.font.bold = True

    p = content_frame.add_paragraph()
    p.text = "注意：系统内力（相互作用力）不改变系统总动量"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.space_before = Pt(10)

def create_application_slide_1(prs):
    """第8页：碰撞问题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "三、动量守恒定律的应用"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "1. 碰撞问题"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = ORANGE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.7))
    content_frame = content_box.text_frame

    collision_types = [
        "弹性碰撞：动量守恒，机械能守恒",
        "非弹性碰撞：动量守恒，机械能不守恒",
        "完全非弹性碰撞：碰撞后结合在一起"
    ]

    for collision_type in collision_types:
        p = content_frame.add_paragraph()
        p.text = f"• {collision_type}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(18)

    p = content_frame.add_paragraph()
    p.text = "特点：作用时间短，内力远大于外力"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True

def create_application_slide_2(prs):
    """第9页：爆炸问题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "2. 爆炸问题"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "特点："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    features = [
        "内力远大于外力，动量守恒",
        "过程极短，位置变化可忽略",
        "其他形式的能转化为动能"
    ]

    for feature in features:
        p = content_frame.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(14)
        p.level = 1

    p = content_frame.add_paragraph()
    p.text = "应用：炮弹发射、炸弹爆炸等"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.space_before = Pt(10)

def create_application_slide_3(prs):
    """第10页：反冲运动"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "3. 反冲运动"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "定义："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "物体在内力作用下分裂为两部分，向相反方向运动"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(24)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "应用实例："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    examples = ["火箭发射", "喷气式飞机", "枪炮后座", "章鱼喷水推进"]

    for example in examples:
        p = content_frame.add_paragraph()
        p.text = f"• {example}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
        p.level = 1

def create_application_slide_4(prs):
    """第11页：人船模型"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "4. 人船模型"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame

    p = content_frame.add_paragraph()
    p.text = "情景："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "人在静止的船上行走，求人和船的位移"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.space_after = Pt(24)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "结论："
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = content_frame.add_paragraph()
    p.text = "m人·s人 = m船·s船"
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = "人和船的位移与质量成反比"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.level = 1

def create_example_slide_1(prs):
    """第12页：例题1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "四、典型例题分析"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 题号
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.5))
    num_frame = num_box.text_frame
    num_frame.text = "例题1"
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(28)
    num_para.font.bold = True
    num_para.font.color.rgb = ORANGE

    # 题目
    question_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(4.9))
    question_frame = question_box.text_frame

    p = question_frame.add_paragraph()
    p.text = "质量为 m₁ = 2kg 的小球 A 以速度 v₁ = 3m/s 与质量为"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

    p = question_frame.add_paragraph()
    p.text = "m₂ = 1kg 的静止小球 B 发生弹性碰撞，求碰撞后两球的速度。"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(30)

    p = question_frame.add_paragraph()
    p.text = "解："
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p = question_frame.add_paragraph()
    p.text = "动量守恒：m₁v₁ = m₁v₁' + m₂v₂'"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "2×3 = 2v₁' + 1v₂'  ... ①"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "机械能守恒：½m₁v₁² = ½m₁v₁'² + ½m₂v₂'²"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "½×2×3² = ½×2×v₁'² + ½×1×v₂'²  ... ②"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "联立①②解得：v₁' = 1 m/s，v₂' = 4 m/s"
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.font.bold = True

def create_example_slide_2(prs):
    """第13页：例题2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "例题2"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 题目
    question_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.4))
    question_frame = question_box.text_frame

    p = question_frame.add_paragraph()
    p.text = "质量为 M 的炮车以仰角 θ 发射一枚质量为 m 的炮弹，"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

    p = question_frame.add_paragraph()
    p.text = "炮弹出口速度为 v₀，求炮车的后退速度。"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(30)

    p = question_frame.add_paragraph()
    p.text = "解：水平方向动量守恒（设炮车后退速度为 V）"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(12)

    p = question_frame.add_paragraph()
    p.text = "MV - m·v₀cosθ = 0"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)

    p = question_frame.add_paragraph()
    p.text = "∴ V = mv₀cosθ / M"
    p.font.size = Pt(26)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.space_after = Pt(18)

    p = question_frame.add_paragraph()
    p.text = "炮车后退速度为 mv₀cosθ / M，方向水平向后"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True

def create_example_slide_3(prs):
    """第14页：例题3（人船模型）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "例题3"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 题目
    question_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.4))
    question_frame = question_box.text_frame

    p = question_frame.add_paragraph()
    p.text = "质量为 m = 60kg 的人站在质量为 M = 180kg 的静止船上，"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

    p = question_frame.add_paragraph()
    p.text = "船长 L = 3m。人从船头走到船尾，求船的位移。"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(30)

    p = question_frame.add_paragraph()
    p.text = "解：根据人船模型"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(12)

    p = question_frame.add_paragraph()
    p.text = "m·s人 = M·s船"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.space_after = Pt(10)

    p = question_frame.add_paragraph()
    p.text = "且 s人 + s船 = L"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.space_after = Pt(10)

    p = question_frame.add_paragraph()
    p.text = "60 × s人 = 180 × s船"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "∴ s船 = L/4 = 3/4 = 0.75 m"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.space_after = Pt(8)

    p = question_frame.add_paragraph()
    p.text = "船向后移动 0.75m"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.font.italic = True

def create_summary_slide(prs):
    """第15页：课堂小结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "课堂小结"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
    content_frame = content_box.text_frame

    points = [
        "动量是描述物体运动状态的物理量，p = mv",
        "动量守恒定律是自然界基本规律之一",
        "适用条件：系统不受外力或合外力为零",
        "典型应用：碰撞、爆炸、反冲运动",
        "解题关键是选择系统和判断守恒条件"
    ]

    for point in points:
        p = content_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.space_after = Pt(18)

def create_closing_slide(prs):
    """第16页：结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 装饰圆圈
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(2.5), Inches(4), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE

    # 文本框
    text_box = slide.shapes.add_textbox(Inches(0), Inches(2.2), Inches(10), Inches(2.6))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "谢谢观看！"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5.2), Inches(6), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "THANK YOU"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

def create_presentation():
    """创建完整的16页PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建所有幻灯片
    create_title_slide(prs)          # 第1页：封面
    create_toc_slide(prs)            # 第2页：目录
    create_concept_slide_1(prs)      # 第3页：动量的定义
    create_concept_slide_2(prs)      # 第4页：动量的单位
    create_concept_slide_3(prs)      # 第5页：动量的变化
    create_law_slide_1(prs)          # 第6页：动量守恒定律内容
    create_law_slide_2(prs)          # 第7页：守恒条件
    create_application_slide_1(prs)  # 第8页：碰撞问题
    create_application_slide_2(prs)  # 第9页：爆炸问题
    create_application_slide_3(prs)  # 第10页：反冲运动
    create_application_slide_4(prs)  # 第11页：人船模型
    create_example_slide_1(prs)      # 第12页：例题1
    create_example_slide_2(prs)      # 第13页：例题2
    create_example_slide_3(prs)      # 第14页：例题3
    create_summary_slide(prs)        # 第15页：课堂小结
    create_closing_slide(prs)        # 第16页：结束页

    # 保存文件
    output_path = "动量守恒定律.pptx"
    prs.save(output_path)
    print(f"PPT已成功生成！")
    print(f"文件路径: {os.path.abspath(output_path)}")
    print(f"总页数: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()
